import sys
import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Myntra Appliqué Design System Color Tokens (http://applique.myntra.com)
    C_BG_PAGE = RGBColor(0xF5, 0xF5, 0xF6)        # #F5F5F6 Light Clean Page
    C_BG_CARD = RGBColor(0xFF, 0xFF, 0xFF)        # #FFFFFF Pure White Card
    C_BG_ACCENT = RGBColor(0xFF, 0xF0, 0xF5)      # #FFF0F5 Soft Pink Card Fill
    C_PRIMARY = RGBColor(0xFF, 0x3F, 0x6C)        # #FF3F6C Myntra Pink
    C_TEXT_PRIMARY = RGBColor(0x28, 0x2C, 0x3F)   # #282C3F Deep Charcoal Slate
    C_TEXT_SECONDARY = RGBColor(0x53, 0x57, 0x66) # #535766 Slate Gray
    C_TEXT_MUTED = RGBColor(0x94, 0x96, 0x9F)     # #94969F Muted Gray
    C_BORDER_LIGHT = RGBColor(0xEA, 0xEA, 0xEC)   # #EAEAEC Subtle Border
    C_GREEN = RGBColor(0x03, 0xA6, 0x85)          # #03A685 Myntra Emerald Green
    C_AMBER = RGBColor(0xD9, 0x77, 0x06)          # #D97706 Warm Gold
    C_DARK_BANNER = RGBColor(0x28, 0x2C, 0x3F)     # #282C3F Header Text

    def set_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_BG_PAGE
        bg.line.fill.background()

    def add_header(slide, slide_num, cat_title, main_headline):
        # Top Header Banner Box
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        header.fill.solid()
        header.fill.fore_color.rgb = C_BG_CARD
        header.line.color.rgb = C_BORDER_LIGHT

        # Left-aligned Section Category Badge (Pill Style)
        tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.15), Inches(3.2), Inches(0.28))
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = C_BG_ACCENT
        tag_box.line.color.rgb = C_PRIMARY

        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = f"SLIDE {slide_num:02d} | {cat_title.upper()}"
        p_tag.font.size = Pt(10)
        p_tag.font.bold = True
        p_tag.font.color.rgb = C_PRIMARY
        p_tag.alignment = PP_ALIGN.CENTER

        # Left-aligned Main Actionable Headline
        tf_head = slide.shapes.add_textbox(Inches(0.6), Inches(0.46), Inches(12.133), Inches(0.6)).text_frame
        tf_head.word_wrap = True
        p_head = tf_head.paragraphs[0]
        p_head.text = main_headline
        p_head.font.size = Pt(19)
        p_head.font.bold = True
        p_head.font.color.rgb = C_TEXT_PRIMARY

    def add_card(slide, left, top, width, height, highlight=False, fill_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        if fill_color:
            shape.fill.fore_color.rgb = fill_color
            shape.line.color.rgb = C_PRIMARY if highlight else C_BORDER_LIGHT
        elif highlight:
            shape.fill.fore_color.rgb = C_BG_ACCENT
            shape.line.color.rgb = C_PRIMARY
        else:
            shape.fill.fore_color.rgb = C_BG_CARD
            shape.line.color.rgb = C_BORDER_LIGHT
        return shape

    # =========================================================================
    # SLIDE 1: CONTEXT & STRATEGIC GOAL (NextLeap Structure)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1)
    add_header(slide1, 1, "Context & Strategic Goal", "Millions Wishlist Fashion Items, but 82% Cool Off Unpurchased after 30 Days")

    # Left Column Card: CURRENT REALITY & WHY IT MATTERS
    add_card(slide1, 0.6, 1.35, 5.9, 5.7, highlight=True)
    tf1_l = slide1.shapes.add_textbox(Inches(0.75), Inches(1.48), Inches(5.6), Inches(5.4)).text_frame
    tf1_l.word_wrap = True
    
    p = tf1_l.paragraphs[0]
    p.text = "CURRENT REALITY (Baseline)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    p_body1 = tf1_l.add_paragraph()
    p_body1.text = "The average fashion shopper accumulates 35–120 wishlisted items over sales and daily browsing. However, 82% of saved products are never purchased within 30 days.\n"
    p_body1.font.size = Pt(11.5)
    p_body1.font.color.rgb = C_TEXT_SECONDARY

    p_h = tf1_l.add_paragraph()
    p_h.text = "WHY IT MATTERS (Hypothesis)"
    p_h.font.size = Pt(14)
    p_h.font.bold = True
    p_h.font.color.rgb = C_AMBER

    p_hbody = tf1_l.add_paragraph()
    p_hbody.text = "Wishlist accumulation without checkout is NOT a lack of interest—it is an early warning signal of sizing hesitation, price paranoia, and list clutter. Users who get trapped in dormant wishlists show lower overall LTV and higher churn risk."
    p_hbody.font.size = Pt(11.5)
    p_hbody.font.color.rgb = C_TEXT_SECONDARY

    # Right Column: STRATEGIC GOAL & 3 CORE OUTCOMES
    add_card(slide1, 6.7, 1.35, 6.0, 5.7)
    tf1_r = slide1.shapes.add_textbox(Inches(6.85), Inches(1.48), Inches(5.7), Inches(5.4)).text_frame
    tf1_r.word_wrap = True

    p = tf1_r.paragraphs[0]
    p.text = "STRATEGIC NORTH STAR GOAL"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_GREEN

    p_g = tf1_r.add_paragraph()
    p_g.text = "Increase 30-Day Wishlist-to-Purchase Conversion Rate by +350 bps (12% baseline → 15.5% target) under strict zero-discount constraints.\n"
    p_g.font.size = Pt(12)
    p_g.font.bold = True
    p_g.font.color.rgb = C_TEXT_PRIMARY

    outcomes = [
        ("Expand Conversion Potential", "Convert high-intent saved items into single-item checkouts without monetary price cuts."),
        ("Eliminate Sizing & Fabric Uncertainty", "Provide social proof try-on reviews & true-to-size ratings directly on wishlist cards."),
        ("Strengthen Long-Term Retain & LTV", "Prevent high-intent saved items from getting buried under 100+ unorganized bookmarks.")
    ]
    for title, desc in outcomes:
        p = tf1_r.add_paragraph()
        p.text = f"✔ {title}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p_d = tf1_r.add_paragraph()
        p_d.text = f"  {desc}"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = C_TEXT_SECONDARY

    # =========================================================================
    # SLIDE 2: HYPOTHESIS & CURRENT DISCOVERY FLOW BREAKDOWN
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2)
    add_header(slide2, 2, "Hypothesis & Breakdown", "Wishlist Engine Fails Because it Optimizes for Bookmarking Rather than Purchase Conversion")

    # Top Hypothesis Banner
    add_card(slide2, 0.6, 1.35, 12.133, 1.25, highlight=True)
    tf2_h = slide2.shapes.add_textbox(Inches(0.75), Inches(1.42), Inches(11.833), Inches(1.1)).text_frame
    tf2_h.word_wrap = True
    p = tf2_h.paragraphs[0]
    p.text = "HYPOTHESIZED STRUCTURAL BREAKDOWN"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    p2 = tf2_h.add_paragraph()
    p2.text = "Myntra's wishlist experience is not broken—it is optimizing for bookmarking volume rather than purchase intent. Shoppers save 50+ items while scrolling, but once saved, the interface offers zero pre-purchase fit validation or occasion organization, causing intent to cool off."
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = C_TEXT_PRIMARY

    # 3 Breakdown Column Cards
    flows = [
        ("1. BOOKMARKING PHASE", "Optimized for Save Volume", "User saves 50–100 items during EORS or commute scrolling.\n\nFlaw:\nNo prompt for size selection or intended occasion."),
        ("2. REVISIT & RETRIEVAL PHASE", "Flat Unfiltered List", "User reopens wishlist 7 days later.\n\nFlaw:\nHigh-intent items get buried under 100+ unorganized cards, causing decision fatigue."),
        ("3. CHECKOUT COMMITMENT PHASE", "Structural Friction Drop-Off", "User considers buying.\n\nFlaw:\nFit anxiety (89.8%), price floor paranoia, and delivery fee halts checkout.")
    ]
    for i, (title, sub, body) in enumerate(flows):
        x = 0.6 + i * 4.1
        add_card(slide2, x, 2.75, 3.93, 4.3)
        tf = slide2.shapes.add_textbox(Inches(x+0.15), Inches(2.88), Inches(3.63), Inches(4.0)).text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY if i==2 else C_TEXT_PRIMARY

        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = C_AMBER

        for line in body.split('\n\n'):
            p_b = tf.add_paragraph()
            p_b.text = line
            p_b.font.size = Pt(11)
            p_b.font.color.rgb = C_TEXT_SECONDARY

    # =========================================================================
    # SLIDE 3: REVIEWLENS - AI DISCOVERY ENGINE (VOC FINDINGS)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3)
    add_header(slide3, 3, "ReviewLens AI Discovery", "ReviewLens Engine Analyzes 20,703 Reviews to Uncover 3 Core Unmet Customer Needs")

    # 3 Unmet Need Cards Across Top
    needs = [
        ("N1. Fit & Fabric Confidence", "2,252 Signals (89.8%)", "Users demand real customer photos on similar Indian body types before purchasing."),
        ("N2. Wishlist Anti-Clutter", "181 Signals (7.2%)", "Users need automatic occasion categorization (Workwear, College, Festive) to find saved items."),
        ("N3. Price Floor Transparency", "91 Signals (3.6%)", "Users want 60-day price floor tracking to buy without fear of post-purchase price drops.")
    ]
    for i, (n_title, n_sig, n_desc) in enumerate(needs):
        x = 0.6 + i * 4.1
        add_card(slide3, x, 1.35, 3.93, 1.8, highlight=(i==0))
        tf = slide3.shapes.add_textbox(Inches(x+0.12), Inches(1.42), Inches(3.69), Inches(1.65)).text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = n_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY if i==0 else C_TEXT_PRIMARY

        p_s = tf.add_paragraph()
        p_s.text = n_sig
        p_s.font.size = Pt(11)
        p_s.font.bold = True
        p_s.font.color.rgb = C_AMBER

        p_d = tf.add_paragraph()
        p_d.text = n_desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = C_TEXT_SECONDARY

    # Bottom Banner Card: KEY TAKEAWAY & CORPUS SOURCES
    add_card(slide3, 0.6, 3.3, 12.133, 3.75)
    tf3_b = slide3.shapes.add_textbox(Inches(0.8), Inches(3.42), Inches(11.733), Inches(3.5)).text_frame
    tf3_b.word_wrap = True

    p = tf3_b.paragraphs[0]
    p.text = "KEY TAKEAWAY FROM REVIEWLENS DISCOVERY ENGINE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    p_t = tf3_b.add_paragraph()
    p_t.text = "Users want Myntra to provide pre-purchase fit reassurance and list curation—NOT extra price cuts. Conversion is blocked by risk aversion, not affordability.\n"
    p_t.font.size = Pt(11.5)
    p_t.font.color.rgb = C_TEXT_PRIMARY

    p_sources = tf3_b.add_paragraph()
    p_sources.text = "Ingested Corpus Breakdown (20,703 Records):"
    p_sources.font.size = Pt(12)
    p_sources.font.bold = True
    p_sources.font.color.rgb = C_GREEN

    sources_list = [
        "• Google Play Store India: 5,000 Android app reviews scraped & keyword-screened.",
        "• Apple App Store (Global & IN): 2,500 iOS storefront reviews processed.",
        "• Reddit r/myntra & r/IndianFashionAddicts: Public discussion threads on sizing and wishlist habits.",
        "• LLM Pipeline: Dual-provider Gemini 1.5 & Groq LLaMA-3.3 automated batch classification."
    ]
    for src in sources_list:
        p_s = tf3_b.add_paragraph()
        p_s.text = src
        p_s.font.size = Pt(11)
        p_s.font.color.rgb = C_TEXT_SECONDARY

    # =========================================================================
    # SLIDE 4: USER RESEARCH / UNDERSTANDING THE USER
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4)
    add_header(slide4, 4, "User Research & Persona", "Surveys (n=30) & Qualitative Interviews (n=6) Validate Sizing & Clutter Blockers")

    # Left Box: Primary Research Summary & Survey Form
    add_card(slide4, 0.6, 1.35, 5.9, 5.7)
    tf4_l = slide4.shapes.add_textbox(Inches(0.75), Inches(1.48), Inches(5.6), Inches(5.4)).text_frame
    tf4_l.word_wrap = True

    p = tf4_l.paragraphs[0]
    p.text = "PRIMARY RESEARCH SUMMARY (n=30 Surveys, n=6 Interviews)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    survey_stats = [
        ("89.8% Sizing Uncertainty", "25 out of 30 respondents stated they hesitate to checkout because brand size charts are inconsistent."),
        ("7.2% Wishlist Disorganization", "22 out of 30 reported their wishlist is an unorganized dump of 50+ items where saved dresses get lost."),
        ("69% Demand Fit Control", "Users want verified real-buyer try-on photos filtered by buyer height & weight.")
    ]
    for st, sd in survey_stats:
        p = tf4_l.add_paragraph()
        p.text = f"• {st}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_PRIMARY
        p_d = tf4_l.add_paragraph()
        p_d.text = f"  {sd}"
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = C_TEXT_SECONDARY

    # Right Box: USER PERSONA CARD
    add_card(slide4, 6.7, 1.35, 6.0, 5.7, highlight=True)
    tf4_r = slide4.shapes.add_textbox(Inches(6.85), Inches(1.48), Inches(5.7), Inches(5.4)).text_frame
    tf4_r.word_wrap = True

    p = tf4_r.paragraphs[0]
    p.text = "TARGET PERSONA: Ananya M. (24, Bangalore)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    p_sub = tf4_r.add_paragraph()
    p_sub.text = "College Student & Frequent Daily Fashion Shopper\n"
    p_sub.font.size = Pt(11)
    p_sub.font.bold = True
    p_sub.font.color.rgb = C_AMBER

    persona_details = [
        ("Goals:", "Discover trendy daily college & weekend outfits within ₹500–₹1,500 budget."),
        ("Needs:", "Personalized recommendations with real-buyer try-on photos matching her 5'4\" height and 52kg weight."),
        ("Pain Points:", "Cannot judge fabric drape from 5'10\" model studio photos; fears getting stuck in return/exchange loops."),
        ("Direct Quote:", "\"I have 35 items saved. I really want to buy the kurti, but every brand's Size M is different. Without seeing it on a real person, I leave it saved for weeks.\"")
    ]
    for k, v in persona_details:
        p = tf4_r.add_paragraph()
        p.text = f"• {k}"
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_PRIMARY
        p_v = tf4_r.add_paragraph()
        p_v.text = f"  {v}"
        p_v.font.size = Pt(10.5)
        p_v.font.italic = (k=="Direct Quote:")
        p_v.font.color.rgb = C_TEXT_SECONDARY

    # =========================================================================
    # SLIDE 5: TARGET SEGMENT & PROBLEM FRAMING CANVAS
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5)
    add_header(slide5, 5, "Target Segment & Problem Canvas", "Focusing on 'High-Intent Accumulators': 10–30 Saved Items Picked by Size Yield 4.2x LTV")

    # 4 Archetype Grid Boxes across left/right
    # Top Left: Who is Target Segment
    add_card(slide5, 0.6, 1.35, 5.9, 2.7, highlight=True)
    tf5_tl = slide5.shapes.add_textbox(Inches(0.75), Inches(1.45), Inches(5.6), Inches(2.4)).text_frame
    tf5_tl.word_wrap = True
    p = tf5_tl.paragraphs[0]
    p.text = "TARGET SEGMENT: High-Intent Accumulators"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    p = tf5_tl.add_paragraph()
    p.text = "Shoppers with 10–30 saved items who have already pre-selected a specific size ('S' or 'M'). They revisit their wishlist 2–3 times weekly and possess high purchasing power but conservative fit risk tolerance."
    p.font.size = Pt(11)
    p.font.color.rgb = C_TEXT_SECONDARY

    # Top Right: Why This Segment
    add_card(slide5, 6.7, 1.35, 6.0, 2.7)
    tf5_tr = slide5.shapes.add_textbox(Inches(6.85), Inches(1.45), Inches(5.7), Inches(2.4)).text_frame
    tf5_tr.word_wrap = True
    p = tf5_tr.paragraphs[0]
    p.text = "WHY THIS SEGMENT? (Highest LTV Potential)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_GREEN
    p = tf5_tr.add_paragraph()
    p.text = "• Strongest Intent Signal: Size already chosen.\n• Highest Conversion Potential: 4.2x higher than casual moodboarders.\n• Solvable Friction: Providing social proof try-ons unlocks immediate checkout."
    p.font.size = Pt(11)
    p.font.color.rgb = C_TEXT_SECONDARY

    # Bottom Full Card: PROBLEM FRAMING CANVAS
    add_card(slide5, 0.6, 4.2, 12.133, 2.85)
    tf5_b = slide5.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.733), Inches(2.65)).text_frame
    tf5_b.word_wrap = True

    p = tf5_b.paragraphs[0]
    p.text = "PROBLEM FRAMING CANVAS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    p_canvas = [
        ("User & Context:", "High-Intent fashion shoppers revisiting their saved wishlist items on Myntra mobile app."),
        ("Trigger:", "Item added during commute or sale browsing; user reopens app 3–7 days later to buy."),
        ("Problem:", "Shopper stops short of checkout due to fit uncertainty and decision fatigue from list clutter."),
        ("Root Cause:", "Studio model photos lack real-body context; zero occasion folder organization; price paranoia."),
        ("Business Impact:", "82% of high-intent saved items remain unmonetized; lower purchase frequency & platform LTV.")
    ]
    for k, v in p_canvas:
        p = tf5_b.add_paragraph()
        p.text = f"• {k} {v}"
        p.font.size = Pt(11)
        p.font.color.rgb = C_TEXT_PRIMARY

    # =========================================================================
    # SLIDE 6: SOLUTION EXPLORATION & RICE PRIORITIZATION
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6)
    add_header(slide6, 6, "Solution Prioritization", "RICE Matrix Ranks FitConfidence Overlay (15.0 Score) as the #1 High-Impact Lever")

    # RICE Table Box
    add_card(slide6, 0.6, 1.35, 12.133, 5.7)
    tf6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.48), Inches(11.733), Inches(5.4)).text_frame
    tf6.word_wrap = True

    p = tf6.paragraphs[0]
    p.text = "RICE PRIORITIZATION MATRIX (Zero-Discount Scope)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    p_sub = tf6.add_paragraph()
    p_sub.text = "Formula: RICE Score = (Reach × Impact × Confidence) / Effort\n"
    p_sub.font.size = Pt(11)
    p_sub.font.color.rgb = C_TEXT_MUTED

    solutions = [
        ("S1: FitConfidence Overlay & Real Buyer Photos", "Addresses 89.8% Fit Anxiety", "R: 9 | I: 10 | C: 9 | E: 6", "15.0", True),
        ("S2: Smart Occasion Folders (Work/College/Festive)", "Addresses 7.2% Wishlist Clutter", "R: 8 | I: 8 | C: 9 | E: 5", "11.5", False),
        ("S3: Smart Value Check (Price When Added Tracking)", "Addresses Price Floor Paranoia", "R: 9 | I: 7 | C: 8 | E: 4", "12.6", False),
        ("S4: Free Delivery Threshold Bundler", "Addresses 3.6% Shipping Fee Friction", "R: 7 | I: 6 | C: 8 | E: 4", "8.4", False)
    ]

    for title, addr, rice_parts, score, is_top in solutions:
        p_t = tf6.add_paragraph()
        p_t.text = f"✔ {title}   [RICE SCORE: {score}]"
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = C_PRIMARY if is_top else C_GREEN

        p_d = tf6.add_paragraph()
        p_d.text = f"   Focus: {addr}   |   Parameters: {rice_parts}"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = C_TEXT_SECONDARY

    p_princ = tf6.add_paragraph()
    p_princ.text = "\nCORE SOLUTION PRINCIPLE:"
    p_princ.font.size = Pt(12)
    p_princ.font.bold = True
    p_princ.font.color.rgb = C_AMBER

    p_princ_body = tf6.add_paragraph()
    p_princ_body.text = "The solution suite must balance automation with light user input—surfacing pre-purchase fit confidence and occasion curation without relying on price cuts or promotional discounts."
    p_princ_body.font.size = Pt(11)
    p_princ_body.font.color.rgb = C_TEXT_SECONDARY

    # =========================================================================
    # SLIDE 7: PRODUCT SOLUTION ANNOUNCEMENT (MYNTRA BINGO CONVERSION SUITE)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7)
    add_header(slide7, 7, "Product Solution Suite", "Announcing 'Myntra Wishlist Conversion Suite': 5 Zero-Discount Product Levers")

    # 3 Value Pillar Cards
    pillars = [
        ("✨ 1. FitConfidence & Buyer Photos", "Social Proof Overlay", "Displays live true-to-size ratings (e.g. '88% True to Size') and 1-tap unedited buyer photos filtered by height (5'4\") & weight (52kg).", True),
        ("📁 2. Smart Occasion Folders", "Anti-Clutter Organization", "Auto-groups 50+ saved items into 'Workwear', 'College', 'Festive', and 'Gym' tabs to eliminate list scroll fatigue.", False),
        ("💡 3. Smart Value Check", "Price Floor & Saved Price", "Shows exact price when item was saved vs current price, proving current price is a real low (zero fake sale stress).", False)
    ]
    for i, (p_title, p_sub, p_desc, is_p1) in enumerate(pillars):
        x = 0.6 + i * 4.1
        add_card(slide7, x, 1.35, 3.93, 5.7, highlight=is_p1)
        tf = slide7.shapes.add_textbox(Inches(x+0.15), Inches(1.48), Inches(3.63), Inches(5.4)).text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY if is_p1 else C_TEXT_PRIMARY

        p_s = tf.add_paragraph()
        p_s.text = p_sub
        p_s.font.size = Pt(11)
        p_s.font.bold = True
        p_s.font.color.rgb = C_AMBER

        p_d = tf.add_paragraph()
        p_d.text = f"\n{p_desc}\n"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = C_TEXT_SECONDARY

        if is_p1:
            p_m = tf.add_paragraph()
            p_m.text = "✔ Solves 89.8% Fit Anxiety"
            p_m.font.size = Pt(11)
            p_m.font.bold = True
            p_m.font.color.rgb = C_GREEN

    # =========================================================================
    # SLIDE 8: WIREFRAMES AND DECODING THE MVP
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8)
    add_header(slide8, 8, "Wireframes & Decoding MVP", "Decoding the Product MVP: 4 Mobile Wireframe Screens & 4 Core AI Transformation Pillars")

    # Left Column: 4 Mobile Wireframe Callout Cards (Grid layout)
    wireframe_screens = [
        ("Screen 1: Login Page", "Introducing Login Page", "Profile chooser (Meera / Rohan) pre-loaded with height/weight attributes to eliminate login friction in testing."),
        ("Screen 2: Home Page", "Home & Wishlist Feed", "Shows all categories, occasion pills ('College', 'Festive') and AI-scored item recommendations."),
        ("Screen 3: Category Suggestion Card", "FitConfidence Overlay", "Displays Basket Fit %, plain-language reason ('Why this?'), and real products in one card without extra navigation."),
        ("Screen 4: Instant Cart & Rescoring", "Tapping Add Updates Cart", "Updates cart instantly and triggers background rescoring; discovery & purchase close in one tap.")
    ]

    for i, (w_title, w_badge, w_desc) in enumerate(wireframe_screens):
        col = i % 2
        row = i // 2
        x = 0.6 + col * 3.0
        y = 1.35 + row * 2.85
        add_card(slide8, x, y, 2.85, 2.7, highlight=(i==2))
        tf = slide8.shapes.add_textbox(Inches(x+0.1), Inches(y+0.1), Inches(2.65), Inches(2.5)).text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = w_title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY

        p_b = tf.add_paragraph()
        p_b.text = f"[{w_badge}]"
        p_b.font.size = Pt(9.5)
        p_b.font.bold = True
        p_b.font.color.rgb = C_GREEN

        p_d = tf.add_paragraph()
        p_d.text = f"\n{w_desc}"
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = C_TEXT_SECONDARY

    # Right Column: 4 Strategic AI & UX Pillars (Reference Image Alignment)
    ai_pillars = [
        ("Why Traditional Systems Fail", "Solves the 'static rule' trap & legacy banners—most systems over-index on reorder history or ignore cart/wishlist context completely."),
        ("What AI Unlocks", "AI treats live basket as real-time signal, reasoning over saved contents on every change instead of a static co-occurrence table."),
        ("How It Transforms UX", "Shifts user from passive reorderer to active discoverer via threshold-gated explainable rails (e.g. 88% fit score on Kurtis)."),
        ("Why AI Is Uniquely Required", "Basket-affinity reasoning ensures relevance & explanation stay personalized in real time, not recalculated from a stale snapshot.")
    ]

    for i, (p_title, p_body) in enumerate(ai_pillars):
        y = 1.35 + i * 1.42
        add_card(slide8, 6.8, y, 5.93, 1.35, highlight=(i==1 or i==2))
        tf = slide8.shapes.add_textbox(Inches(6.95), Inches(y+0.08), Inches(5.63), Inches(1.18)).text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY if (i==1 or i==2) else C_TEXT_PRIMARY

        p_b = tf.add_paragraph()
        p_b.text = p_body
        p_b.font.size = Pt(9.5)
        p_b.font.color.rgb = C_TEXT_SECONDARY

    # =========================================================================
    # SLIDE 9: MEASURING SUCCESS / RISKS & MITIGATIONS
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_bg(slide9)
    add_header(slide9, 9, "Success Metrics & Risks", "Tracking 30-Day Wishlist Conversion (+350 bps) Against Card-to-Bag CTR (+18%)")

    # 3 Column Metric Cards: North Star / L1 & L2 / Guardrails & Risks
    cols = [
        ("NORTH STAR METRIC", "Lagging Metric (30 Days)", [
            ("30-Day Wishlist Conversion Rate", "Target: +350 bps (12% → 15.5%)", "Percentage of unique users purchasing at least one wishlisted item within 30 days.")
        ], C_PRIMARY),
        ("L1 & L2 METRICS", "Leading Indicators (0–14 Days)", [
            ("Wishlist Card-to-Bag CTR", "Target: +18% lift", "Measures size selection into bag."),
            ("Buyer Photo Modal Adoption", "Target: 35% usage", "Tracks photo reel engagement."),
            ("Occasion Folder Reopen Rate", "Target: +25% reopen", "Measures anti-clutter impact.")
        ], C_GREEN),
        ("RISKS & GUARDRAILS", "Risk Controls & Latency", [
            ("Return & Exchange Rate", "Guardrail: Must NOT rise >0.5%", "Confirms sizing accuracy."),
            ("Photo Moderation Queue", "Mitigation: AI NSFW filter", "Prevents invalid buyer photos."),
            ("CDN Model Latency", "Guardrail: <50ms overhead", "Pre-computes price sparklines.")
        ], C_AMBER)
    ]

    for i, (c_title, c_sub, c_items, c_color) in enumerate(cols):
        x = 0.6 + i * 4.1
        add_card(slide9, x, 1.35, 3.93, 5.7, highlight=(i==0))
        tf = slide9.shapes.add_textbox(Inches(x+0.15), Inches(1.48), Inches(3.63), Inches(5.4)).text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = c_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = c_color

        p_s = tf.add_paragraph()
        p_s.text = c_sub
        p_s.font.size = Pt(11)
        p_s.font.bold = True
        p_s.font.color.rgb = C_TEXT_MUTED

        for m_name, m_tgt, m_desc in c_items:
            p_n = tf.add_paragraph()
            p_n.text = f"\n• {m_name}"
            p_n.font.size = Pt(11.5)
            p_n.font.bold = True
            p_n.font.color.rgb = C_TEXT_PRIMARY

            p_t = tf.add_paragraph()
            p_t.text = f"  [{m_tgt}]"
            p_t.font.size = Pt(10.5)
            p_t.font.bold = True
            p_t.font.color.rgb = c_color

            p_d = tf.add_paragraph()
            p_d.text = f"  {m_desc}"
            p_d.font.size = Pt(10)
            p_d.font.color.rgb = C_TEXT_SECONDARY

    # =========================================================================
    # SLIDE 10: ROADMAP: PHASE-WISE ROLLOUT & NEXT STEPS
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_bg(slide10)
    add_header(slide10, 10, "Rollout Roadmap & Next Steps", "Phased 4-Stage Rollout Strategy Maximizes Zero-Discount LTV Impact across Myntra App")

    # 4 Phase Cards across the slide width
    phases_s10 = [
        ("Phase 1: Alpha & Core Loop", "Week 1–3 | 1% MAU", "Invite-only High-Intent Wishlist Accumulators; validate FitConfidence scoring (88% rating) & height/weight buyer photo reels over 2-week window.", True),
        ("Phase 2: Targeted Beta", "Week 4–6 | 10% MAU", "Expand to Metro + Tier-2/3 fashion clusters; test Smart Occasion Folders, dismiss fatigue & CDN load latency (<50ms).", False),
        ("Phase 3: General Availability", "Week 7+ | 100% Rollout", "Full rollout across iOS & Android app. FitConfidence live by default on Wishlist & Checkout drawer; zero-discount gate.", False),
        ("Phase 4: Non-Monetary Layer", "Quarterly Follow-up", "Introduce restock alerts for wishlisted sizes, cross-category occasion lookbook pairings, and per-category risk weighting.", False)
    ]

    for i, (ph_title, ph_meta, ph_desc, is_highlight) in enumerate(phases_s10):
        x = 0.6 + i * 3.05
        add_card(slide10, x, 1.35, 2.9, 5.7, highlight=is_highlight)
        tf = slide10.shapes.add_textbox(Inches(x+0.1), Inches(1.48), Inches(2.7), Inches(5.4)).text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"STAGE 0{i+1}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY if is_highlight else C_AMBER

        p_t = tf.add_paragraph()
        p_t.text = ph_title
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = C_TEXT_PRIMARY

        p_m = tf.add_paragraph()
        p_m.text = f"[{ph_meta}]"
        p_m.font.size = Pt(10)
        p_m.font.bold = True
        p_m.font.color.rgb = C_GREEN

        p_d = tf.add_paragraph()
        p_d.text = f"\n{ph_desc}"
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = C_TEXT_SECONDARY

    # =========================================================================
    # SLIDE 11: TIERS OF MONETIZATION, DISTRIBUTION CHANNELS & RISKS
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_bg(slide11)
    add_header(slide11, 11, "Monetization, Channels & Risks", "Tiers of Monetization, In-App Distribution Channels & Risk Mitigation Architecture")

    # Col 1: Tiers of Monetization
    add_card(slide11, 0.6, 1.35, 3.8, 5.7, highlight=True)
    tf11_1 = slide11.shapes.add_textbox(Inches(0.75), Inches(1.48), Inches(3.5), Inches(5.4)).text_frame
    tf11_1.word_wrap = True

    p = tf11_1.paragraphs[0]
    p.text = "TIERS OF MONETIZATION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    m_tiers = [
        ("Free Tier (Organic Discovery)", "AI-scored FitConfidence overlays shown at no cost; Myntra absorbs compute cost as an LTV investment."),
        ("Brand Partner Co-Op (B2B)", "Category fashion brands pay for priority inclusion in candidate pool — still strictly gated by 50% relevance & 88% fit score."),
        ("Myntra Insider Loyalty", "Loyalty members get first access to verified buyer photo reels, 1-tap occasion customization & free doorstep return validation.")
    ]
    for t_name, t_desc in m_tiers:
        p = tf11_1.add_paragraph()
        p.text = f"\n• {t_name}"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_PRIMARY
        p_d = tf11_1.add_paragraph()
        p_d.text = f"  {t_desc}"
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = C_TEXT_SECONDARY

    # Col 2: Distribution Channels
    add_card(slide11, 4.7, 1.35, 3.8, 5.7)
    tf11_2 = slide11.shapes.add_textbox(Inches(4.85), Inches(1.48), Inches(3.5), Inches(5.4)).text_frame
    tf11_2.word_wrap = True

    p = tf11_2.paragraphs[0]
    p.text = "DISTRIBUTION CHANNELS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_AMBER

    channels = [
        ("Internal Triggers", "Dismiss-trap intercept: 2 dismissals triggers quick 'What would help size choice?' prompt. Cart-page live widget: FitConfidence rail visible directly on Wishlist/Cart."),
        ("External Triggers", "Price Floor & Stock Nudge: WhatsApp/SMS trigger when wishlisted item reaches 60-day price floor. Notification re-engagement: Push nudge if occasion folder viewed 48h ago.")
    ]
    for c_name, c_desc in channels:
        p = tf11_2.add_paragraph()
        p.text = f"\n📲 {c_name}"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_PRIMARY
        p_d = tf11_2.add_paragraph()
        p_d.text = f"  {c_desc}"
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = C_TEXT_SECONDARY

    # Col 3: Risks & Mitigations
    add_card(slide11, 8.8, 1.35, 3.93, 5.7)
    tf11_3 = slide11.shapes.add_textbox(Inches(8.95), Inches(1.48), Inches(3.63), Inches(5.4)).text_frame
    tf11_3.word_wrap = True

    p = tf11_3.paragraphs[0]
    p.text = "RISKS & MITIGATIONS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_GREEN

    risks = [
        ("Technical: Latency Drag", "Instant Heuristic Fallback: If AI call doesn't return within ~1.5s, fall back to static brand fit rating — checkout never waits."),
        ("Engagement: Nudge Fatigue", "Frequency Cap + Mute: Cap suggestions per category per week; let users mute category type outright rather than dismiss every session."),
        ("Analytical: Conflicting Signals", "Signal Decay: Weight explicit negative feedback down over time so stale 'Irrelevant' flag doesn't permanently suppress active categories.")
    ]
    for r_name, r_desc in risks:
        p = tf11_3.add_paragraph()
        p.text = f"\n⚠️ {r_name}"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_PRIMARY
        p_d = tf11_3.add_paragraph()
        p_d.text = f"  {r_desc}"
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = C_TEXT_SECONDARY

    # Save Presentation
    output_filename = "myntra-wishlist-conversion-final.pptx"
    prs.save(output_filename)
    print(f"SUCCESS: Created NextLeap-structured presentation '{output_filename}' with {len(prs.slides)} slides.")
    
    try:
        prs.save("myntra-wishlist-conversion.pptx")
        print("SUCCESS: Also updated 'myntra-wishlist-conversion.pptx'.")
    except Exception as e:
        print("Note: 'myntra-wishlist-conversion.pptx' is locked by PowerPoint; saved as 'myntra-wishlist-conversion-final.pptx'.")

if __name__ == "__main__":
    create_deck()

